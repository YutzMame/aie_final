import json
import os
import boto3
import io
import re
import traceback
import uuid
from pptx import Presentation
import urllib.parse
from datetime import datetime

# --- AWSクライアントの初期化 ---
MODEL_ID = os.environ.get("MODEL_ID", "us.amazon.nova-lite-v1:0")
AWS_REGION = os.environ.get("AWS_REGION", "us-east-1")
TABLE_NAME = os.environ.get("TABLE_NAME")

bedrock_runtime = boto3.client(service_name="bedrock-runtime", region_name=AWS_REGION)
s3_client = boto3.client("s3")
dynamodb = boto3.resource("dynamodb")
table = dynamodb.Table(TABLE_NAME)


def generate_qa_from_text(lecture_text, num_questions, difficulty):
    """
    抽出されたテキストを元に、Bedrockを呼び出してQAを生成する共通関数。
    """
    print(
        f"Generating {num_questions} QAs with difficulty '{difficulty}' from extracted text."
    )

    system_prompt = f"""
あなたは、講義内容から学習者の理解度を測るための問題を作成する専門家です。
以下のルールに従って、与えられた講義内容から質の高いQAセットを作成してください。
# ルール
- 質問形式は「一択選択式」「記述式」をバランス良く含めること。ただし、「記述式」の問題数は少なくすること。
- {num_questions}個の問題を、難易度「{difficulty}」で作成すること。
- 回答には、なぜそれが正解なのかの短い解説を必ず含めること。
- 出力は必ず指定されたJSON形式のみとし、前後に余計な文章は含めないこと。
# JSON形式
{{
  "qa_set": [
    {{
      "question_id": 1,
      "difficulty": "易",
      "type": "一択選択式",
      "question_text": "質問文",
      "options": ["選択肢A", "選択肢B", "選択肢C", "選択肢D"],
      "answer": "正解の選択肢",
      "explanation": "解説文"
    }}
  ]
}}
"""
    user_prompt = f"--- 講義内容 ---\n{lecture_text}"

    # Novaモデルが期待する正しいmessages形式のリクエストボディ
    request_body = {
        "schemaVersion": "messages-v1",
        "system": [{"text": system_prompt}],
        "messages": [{"role": "user", "content": [{"text": user_prompt}]}],
        "inferenceConfig": {
            "maxTokens": 4096,
            "stopSequences": [],
            "temperature": 0.7,
            "topP": 0.9,
        },
    }

    print(f"Calling Bedrock with payload...")
    response = bedrock_runtime.invoke_model(
        body=json.dumps(request_body), modelId=MODEL_ID
    )

    response_body = json.loads(response.get("body").read())
    qa_result_text = (
        response_body.get("output", {})
        .get("message", {})
        .get("content", [{}])[0]
        .get("text")
    )

    if not qa_result_text:
        raise Exception("モデルの応答からテキストを抽出できませんでした。")

    json_match = re.search(r"\{.*\}", qa_result_text, re.DOTALL)
    if json_match:
        return json.loads(json_match.group(0))
    else:
        raise Exception("モデルの応答から有効なJSONを抽出できませんでした。")


def handler(event, context):
    print(f"Received S3 event: {json.dumps(event)}")

    # S3イベントからバケット名とオブジェクトキーを取得
    bucket = event["Records"][0]["s3"]["bucket"]["name"]
    key = urllib.parse.unquote_plus(
        event["Records"][0]["s3"]["object"]["key"], encoding="utf-8"
    )

    try:
        # S3オブジェクトのメタデータを取得
        s3_object = s3_client.head_object(Bucket=bucket, Key=key)
        metadata = s3_object.get("Metadata", {})
        theme = metadata.get("theme", "untitled")
        lecture_number = int(metadata.get("lecture_number", 1))
        num_questions = int(metadata.get("num_questions", 5))
        difficulty = metadata.get("difficulty", "中")

        # S3からファイルをダウンロード
        response = s3_client.get_object(Bucket=bucket, Key=key)
        ppt_content = response["Body"].read()

        # メモリ上でファイルを扱う
        ppt_stream = io.BytesIO(ppt_content)
        presentation = Presentation(ppt_stream)

        # 全スライドからテキストを抽出
        extracted_text = "\n".join(
            shape.text
            for slide in presentation.slides
            for shape in slide.shapes
            if hasattr(shape, "text")
        )
        print(f"Extracted {len(extracted_text)} characters from {key}.")

        # 抽出したテキストを使ってQAを生成
        qa_json = generate_qa_from_text(extracted_text, num_questions, difficulty)

        # DynamoDBに保存
        qa_set_id = str(uuid.uuid4())
        item_to_save = {
            "qa_set_id": qa_set_id,
            "qa_data": qa_json,
            "theme": theme,
            "lecture_number": lecture_number,
            "source_file": key,
            "created_at": datetime.utcnow().isoformat(),
        }
        table.put_item(Item=item_to_save)

        print(f"Successfully processed {key} and saved to DynamoDB.")
        return {"status": "success"}

    except Exception as e:
        print(f"ERROR processing file {key} from bucket {bucket}.")
        print(traceback.format_exc())
        raise e
